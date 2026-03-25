"""
redesign_pptx.py
================
Reads an existing .pptx file, extracts ALL text content from every slide,
and produces a brand-new, visually redesigned presentation with a
corporate/professional blue theme while keeping the original text intact.

Usage:
    python redesign_pptx.py [input_file] [output_file]

Defaults:
    input_file  = presentacion_final_IA.pptx
    output_file = presentacion_final_IA_redesigned.pptx

Configuration constants (see section below) let you adjust colors, fonts,
and sizes without touching the rest of the code.
"""

import sys
import copy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# CONFIGURATION — tweak colors, fonts and sizes here
# ---------------------------------------------------------------------------

# Input / output filenames (can be overridden via CLI args)
INPUT_FILE = "presentacion_final_IA.pptx"
OUTPUT_FILE = "presentacion_final_IA_redesigned.pptx"

# Slide dimensions — widescreen 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# --- Color palette ---
C_NAVY       = RGBColor(0x1B, 0x2A, 0x4A)   # dark navy — title/closing bg
C_CORP_BLUE  = RGBColor(0x2E, 0x50, 0x90)   # corporate blue — section bg / header bar
C_ACCENT     = RGBColor(0x4A, 0x90, 0xD9)   # light blue — accent line / side bar
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)   # white — text on dark backgrounds
C_DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)   # dark gray — body text on light bg
C_LIGHT_BG   = RGBColor(0xF0, 0xF2, 0xF5)   # very light gray — content slide bg
C_FOOTER_BG  = RGBColor(0xE8, 0xEC, 0xF2)   # footer bar background on content slides

# --- Typography ---
FONT_TITLE   = "Calibri"
FONT_BODY    = "Calibri"

PT_COVER_TITLE   = Pt(36)
PT_COVER_SUB     = Pt(20)
PT_SECTION_TITLE = Pt(32)
PT_CONTENT_TITLE = Pt(22)
PT_BODY          = Pt(16)
PT_BULLET_L0     = Pt(16)
PT_BULLET_L1     = Pt(14)
PT_FOOTER        = Pt(10)
PT_SLIDE_NUM     = Pt(10)

# --- Layout measurements (in inches) ---
HEADER_BAR_H  = Inches(1.05)   # height of the blue header bar on content slides
FOOTER_BAR_H  = Inches(0.30)   # height of footer bar
SIDE_BAR_W    = Inches(0.08)   # width of left accent bar
ACCENT_LINE_H = Inches(0.05)   # height of decorative accent line

# Margin inside text boxes
MARGIN_L = Inches(0.15)

# ---------------------------------------------------------------------------
# HELPERS — low-level drawing primitives
# ---------------------------------------------------------------------------

def _set_bg_color(slide, rgb: RGBColor):
    """Fill an entire slide background with a solid RGB color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_rect(slide, left, top, width, height, fill_rgb: RGBColor, line_rgb=None):
    """Add a filled rectangle shape to a slide. Returns the shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()   # no border
    else:
        shape.line.color.rgb = line_rgb
    return shape


def _add_text_box(slide, left, top, width, height, text, font_name, font_size,
                  bold=False, italic=False, color=C_DARK_TEXT,
                  align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box with a single paragraph. Returns the text frame."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    # Remove default margins so text sits flush inside the box
    tf.margin_left   = Pt(0)
    tf.margin_right  = Pt(0)
    tf.margin_top    = Pt(2)
    tf.margin_bottom = Pt(2)

    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def _set_para_spacing(para, space_before_pt=4, space_after_pt=4):
    """Set spacing before/after a paragraph (in points)."""
    pPr = para._pPr
    if pPr is None:
        pPr = para._p.get_or_add_pPr()
    # Use the lxml API directly to append spacing child elements
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
    spcPts.set("val", str(int(space_before_pt * 100)))

    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    spcPts2 = etree.SubElement(spcAft, qn("a:spcPts"))
    spcPts2.set("val", str(int(space_after_pt * 100)))


# ---------------------------------------------------------------------------
# TEXT EXTRACTION
# ---------------------------------------------------------------------------

def _extract_shape_text(shape):
    """
    Recursively extract text paragraphs from a shape.
    Returns a list of dicts: {text, level, is_title}.
    """
    results = []

    # Grouped shapes — recurse into each child
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            results.extend(_extract_shape_text(child))
        return results

    # Tables — extract cell by cell
    if shape.has_table:
        for row in shape.table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text_frame.text.strip()
                if cell_text:
                    row_texts.append(cell_text)
            if row_texts:
                results.append({"text": " | ".join(row_texts), "level": 0, "is_title": False})
        return results

    # Normal text frames
    if shape.has_text_frame:
        is_title = "title" in shape.name.lower()
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                results.append({
                    "text": text,
                    "level": para.level,
                    "is_title": is_title,
                })

    return results


def extract_slide_data(prs: Presentation):
    """
    Walk every slide in *prs* and return a list of slide-data dicts:
    {
        'index': int,               # 0-based
        'paragraphs': [             # ordered list of extracted text paragraphs
            {'text': str, 'level': int, 'is_title': bool},
            ...
        ],
        'notes': str,               # speaker notes (may be empty)
    }
    """
    slides_data = []
    for idx, slide in enumerate(prs.slides):
        paras = []
        for shape in slide.shapes:
            paras.extend(_extract_shape_text(shape))

        # Speaker notes
        notes_text = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                notes_text = notes_tf.text.strip()

        slides_data.append({
            "index": idx,
            "paragraphs": paras,
            "notes": notes_text,
        })

    return slides_data


# ---------------------------------------------------------------------------
# SLIDE-TYPE DETECTION
# ---------------------------------------------------------------------------

def detect_slide_type(slide_data: dict, total_slides: int) -> str:
    """
    Classify a slide as one of:
      'cover'    — first slide
      'closing'  — last slide (if it looks like a closing/thank-you slide)
      'section'  — section divider (few text lines, no deep bullets)
      'content'  — regular content slide
    """
    idx = slide_data["index"]
    paras = slide_data["paragraphs"]

    if idx == 0:
        return "cover"

    non_empty = [p for p in paras if p["text"]]

    if idx == total_slides - 1:
        # Closing slide heuristic: few paragraphs or keywords like 'gracias', 'conclus'
        lower_texts = " ".join(p["text"].lower() for p in non_empty)
        closing_keywords = ("gracias", "thank", "fin", "conclusion", "conclusi", "cierre")
        if len(non_empty) <= 3 or any(kw in lower_texts for kw in closing_keywords):
            return "closing"

    # Section slide: only a title (or title + 1 short line), no level-1 bullets
    has_bullets = any(p["level"] > 0 for p in non_empty)
    body_paras = [p for p in non_empty if not p["is_title"]]
    if not has_bullets and len(body_paras) <= 1:
        return "section"

    return "content"


# ---------------------------------------------------------------------------
# SLIDE BUILDERS
# ---------------------------------------------------------------------------

def _common_footer(slide, slide_num: int):
    """Add a subtle footer bar with slide number at the bottom of any slide."""
    footer_top = SLIDE_H - FOOTER_BAR_H
    _add_rect(slide, Inches(0), footer_top, SLIDE_W, FOOTER_BAR_H, C_FOOTER_BG)
    # Slide number (right-aligned)
    _add_text_box(
        slide,
        left=SLIDE_W - Inches(1.2),
        top=footer_top + Inches(0.04),
        width=Inches(1.0),
        height=FOOTER_BAR_H - Inches(0.08),
        text=str(slide_num),
        font_name=FONT_BODY,
        font_size=PT_SLIDE_NUM,
        color=C_CORP_BLUE,
        align=PP_ALIGN.RIGHT,
    )


def build_cover_slide(slide, slide_data: dict):
    """
    Cover/title slide:
      - Dark navy background
      - Large centered title
      - Horizontal accent line beneath title
      - Subtitle / author info below
    """
    _set_bg_color(slide, C_NAVY)

    # Decorative top stripe
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), C_CORP_BLUE)

    paras = slide_data["paragraphs"]
    title_paras  = [p for p in paras if p["is_title"]]
    body_paras   = [p for p in paras if not p["is_title"]]

    title_text = title_paras[0]["text"] if title_paras else ""

    # ---- Title text box ----
    title_top  = Inches(1.6)
    title_h    = Inches(2.2)
    title_left = Inches(1.2)
    title_w    = SLIDE_W - Inches(2.4)

    txb = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf  = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.name  = FONT_TITLE
    run.font.size  = PT_COVER_TITLE
    run.font.bold  = True
    run.font.color.rgb = C_WHITE

    # ---- Accent line ----
    accent_top = title_top + title_h + Inches(0.15)
    _add_rect(slide,
              Inches(2.5), accent_top,
              SLIDE_W - Inches(5.0), ACCENT_LINE_H,
              C_ACCENT)

    # ---- Subtitle / body ----
    sub_top = accent_top + ACCENT_LINE_H + Inches(0.35)
    sub_h   = SLIDE_H - sub_top - Inches(0.6)

    txb2 = slide.shapes.add_textbox(title_left, sub_top, title_w, sub_h)
    tf2  = txb2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = Pt(0)
    tf2.margin_top  = tf2.margin_bottom = Pt(2)

    first = True
    for p_data in body_paras:
        if first:
            para = tf2.paragraphs[0]
            first = False
        else:
            para = tf2.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = p_data["text"]
        run.font.name  = FONT_BODY
        run.font.size  = PT_COVER_SUB
        run.font.color.rgb = RGBColor(0xC0, 0xD0, 0xE8)   # soft blue-white

    # Decorative bottom stripe
    _add_rect(slide, Inches(0), SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), C_CORP_BLUE)


def build_section_slide(slide, slide_data: dict):
    """
    Section divider slide:
      - Corporate blue background
      - Centered title in white
      - Thin accent line above/below title
      - Optional one-liner subtitle
    """
    _set_bg_color(slide, C_CORP_BLUE)

    paras = slide_data["paragraphs"]
    title_paras = [p for p in paras if p["is_title"]]
    body_paras  = [p for p in paras if not p["is_title"]]

    title_text = title_paras[0]["text"] if title_paras else ""
    sub_text   = body_paras[0]["text"] if body_paras else ""

    center_y = SLIDE_H / 2

    # Accent line above
    _add_rect(slide,
              Inches(2.0), center_y - Inches(1.4),
              SLIDE_W - Inches(4.0), ACCENT_LINE_H,
              C_ACCENT)

    # Title
    t_top = center_y - Inches(1.2)
    t_h   = Inches(1.4)
    txb = slide.shapes.add_textbox(Inches(1.0), t_top, SLIDE_W - Inches(2.0), t_h)
    tf  = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top  = tf.margin_bottom = Pt(0)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    run.font.name  = FONT_TITLE
    run.font.size  = PT_SECTION_TITLE
    run.font.bold  = True
    run.font.color.rgb = C_WHITE

    # Accent line below
    line_top = t_top + t_h + Inches(0.12)
    _add_rect(slide,
              Inches(2.0), line_top,
              SLIDE_W - Inches(4.0), ACCENT_LINE_H,
              C_ACCENT)

    # Optional subtitle
    if sub_text:
        txb2 = slide.shapes.add_textbox(
            Inches(1.5), line_top + Inches(0.3),
            SLIDE_W - Inches(3.0), Inches(0.8)
        )
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_right = tf2.margin_top = tf2.margin_bottom = Pt(0)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = sub_text
        run2.font.name  = FONT_BODY
        run2.font.size  = PT_COVER_SUB
        run2.font.color.rgb = RGBColor(0xC0, 0xD0, 0xE8)


def build_closing_slide(slide, slide_data: dict):
    """
    Closing/thank-you slide:
      - Dark navy background (same as cover)
      - Centered content
    """
    # Reuse the cover style for the closing slide
    build_cover_slide(slide, slide_data)


def build_content_slide(slide, slide_data: dict, slide_num: int):
    """
    Standard content slide:
      - White/light-gray background
      - Dark blue header bar at the top with the slide title in white
      - Thin accent line under the header bar
      - Thin vertical accent bar on the left edge of the content area
      - Body text with proper bullet hierarchy
      - Footer bar with slide number
    """
    # Background
    _set_bg_color(slide, C_LIGHT_BG)

    paras = slide_data["paragraphs"]
    title_paras = [p for p in paras if p["is_title"]]
    body_paras  = [p for p in paras if not p["is_title"]]

    title_text = title_paras[0]["text"] if title_paras else ""

    # ---- Header bar ----
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_BAR_H, C_CORP_BLUE)

    # Title text inside header bar
    title_left = Inches(0.25)
    title_top  = Inches(0.08)
    title_w    = SLIDE_W - Inches(0.5)
    title_h    = HEADER_BAR_H - Inches(0.16)

    txb = slide.shapes.add_textbox(title_left, title_top, title_w, title_h)
    tf  = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top  = tf.margin_bottom = Pt(2)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name  = FONT_TITLE
    run.font.size  = PT_CONTENT_TITLE
    run.font.bold  = True
    run.font.color.rgb = C_WHITE

    # ---- Accent line under header ----
    _add_rect(slide,
              Inches(0), HEADER_BAR_H,
              SLIDE_W, ACCENT_LINE_H,
              C_ACCENT)

    # ---- Left accent bar ----
    content_top  = HEADER_BAR_H + ACCENT_LINE_H + Inches(0.12)
    content_h    = SLIDE_H - content_top - FOOTER_BAR_H - Inches(0.12)
    _add_rect(slide,
              Inches(0), content_top,
              SIDE_BAR_W, content_h,
              C_ACCENT)

    # ---- Body text ----
    body_left = SIDE_BAR_W + Inches(0.22)
    body_top  = content_top + Inches(0.08)
    body_w    = SLIDE_W - body_left - Inches(0.22)
    body_h    = content_h - Inches(0.16)

    if body_paras:
        txb2 = slide.shapes.add_textbox(body_left, body_top, body_w, body_h)
        tf2  = txb2.text_frame
        tf2.word_wrap = True
        tf2.margin_left  = MARGIN_L
        tf2.margin_right = Inches(0.1)
        tf2.margin_top   = Pt(2)
        tf2.margin_bottom = Pt(2)

        first = True
        for p_data in body_paras:
            if first:
                para = tf2.paragraphs[0]
                first = False
            else:
                para = tf2.add_paragraph()

            level = p_data["level"]
            text  = p_data["text"]

            # Indentation by level
            indent_inches = level * 0.25
            para.level = min(level, 8)  # pptx supports levels 0-8

            run = para.add_run()

            # Bullet character prefix for visual clarity
            if level == 0:
                run.text = text
                run.font.size  = PT_BULLET_L0
                run.font.bold  = False
                run.font.color.rgb = C_DARK_TEXT
            else:
                run.text = text
                run.font.size  = PT_BULLET_L1
                run.font.bold  = False
                run.font.color.rgb = RGBColor(0x44, 0x44, 0x55)

            run.font.name = FONT_BODY

    # ---- Footer ----
    _common_footer(slide, slide_num)


# ---------------------------------------------------------------------------
# MAIN BUILDER
# ---------------------------------------------------------------------------

def redesign(input_path: str, output_path: str):
    """Read *input_path*, redesign, and write to *output_path*."""

    print(f"Reading: {input_path}")
    src_prs = Presentation(input_path)

    # Extract all slide text before building the new presentation
    slides_data = extract_slide_data(src_prs)
    total = len(slides_data)
    print(f"  Found {total} slide(s)")

    # Create a blank presentation at the correct widescreen size
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use a completely blank slide layout (index 6 in the default theme)
    blank_layout = prs.slide_layouts[6]

    slide_num = 0  # counter for footer numbers (only incremented on content slides)

    for sd in slides_data:
        idx   = sd["index"]
        stype = detect_slide_type(sd, total)

        slide = prs.slides.add_slide(blank_layout)

        print(f"  Slide {idx+1:2d} [{stype:8s}]  "
              f"{sd['paragraphs'][0]['text'][:55] if sd['paragraphs'] else '(empty)'}")

        if stype == "cover":
            build_cover_slide(slide, sd)

        elif stype == "section":
            build_section_slide(slide, sd)

        elif stype == "closing":
            build_closing_slide(slide, sd)

        else:  # content
            slide_num += 1
            build_content_slide(slide, sd, slide_num)

    prs.save(output_path)
    print(f"\nSaved redesigned presentation → {output_path}")


# ---------------------------------------------------------------------------
# ENTRY POINT
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    in_file  = sys.argv[1] if len(sys.argv) > 1 else INPUT_FILE
    out_file = sys.argv[2] if len(sys.argv) > 2 else OUTPUT_FILE
    redesign(in_file, out_file)
